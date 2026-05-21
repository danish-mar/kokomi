import os
import json
import uuid
import time
import asyncio
import datetime
import contextvars
from typing import Dict, Any, List, Optional, TypedDict
from pydantic import BaseModel, Field

# Storage context for active workflow
active_storage_dir = contextvars.ContextVar("active_storage_dir", default="")
active_workflow_title = contextvars.ContextVar("active_workflow_title", default="")

# LangChain / LangGraph imports
from langchain_core.messages import SystemMessage, HumanMessage, ToolMessage, AIMessage
from langchain_core.tools import tool
from langgraph.graph import StateGraph, END

# App imports
from app.config import DATA_DIR, GROQ_API_KEY, GOOGLE_API_KEY
from app.storage import _load, _save, load_prefs
from app.llm import get_llm, get_atlas_llm
from app.memory import search_memories

# PDF creation library
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, HRFlowable, Preformatted
from reportlab.lib import colors
from reportlab.lib.units import inch
import re as _re

# ── Storage File for Workflows ───────────────────────────────────────
WORKFLOWS_FILE = os.path.join(DATA_DIR, "multi_agent_workflows.json")

_WORKFLOWS_CACHE: dict = {}

def load_workflows() -> dict:
    global _WORKFLOWS_CACHE
    if not _WORKFLOWS_CACHE:
        if not os.path.exists(WORKFLOWS_FILE):
            _WORKFLOWS_CACHE = {}
        else:
            data = _load(WORKFLOWS_FILE)
            if isinstance(data, dict):
                _WORKFLOWS_CACHE = data
            else:
                _WORKFLOWS_CACHE = {}
    return _WORKFLOWS_CACHE

def save_workflows(data: dict) -> None:
    global _WORKFLOWS_CACHE
    _WORKFLOWS_CACHE = data
    try:
        loop = asyncio.get_running_loop()
        # Persist to disk asynchronously in a background thread to prevent blocking the event loop
        loop.create_task(asyncio.to_thread(_save, WORKFLOWS_FILE, data))
    except RuntimeError:
        # Fallback to synchronous save if not running inside an event loop
        _save(WORKFLOWS_FILE, data)

def clean_markdown_if_json(content: str) -> str:
    if not isinstance(content, str):
        return str(content)
    content_stripped = content.strip()
    if content_stripped.startswith('{') and content_stripped.endswith('}'):
        try:
            import json as _json_clean
            data = _json_clean.loads(content_stripped)
            if isinstance(data, dict):
                for key in ["markdown", "result", "content", "text", "body"]:
                    if data.get(key) and isinstance(data[key], str) and len(data[key].strip()) > 50:
                        return data[key]
                longest_str = ""
                for k, v in data.items():
                    if isinstance(v, str) and len(v) > len(longest_str):
                        longest_str = v
                if len(longest_str.strip()) > 50:
                    return longest_str
        except Exception:
            pass
    return content


# ── Typed Schemas ────────────────────────────────────────────────────

class TaskDict(TypedDict, total=False):
    task_id: str
    title: str
    description: str
    worker_type: str
    depends_on: List[str]
    allowed_tools: List[str]
    success_criteria: str
    expected_output_schema: Dict[str, Any]
    status: str  # "pending", "running", "completed", "retrying", "failed_transient", "failed_validation", "blocked", "failed_terminal", "aborted"
    output: Any
    artifacts: List[str]
    retries: int
    error: Optional[str]
    timestamps: Dict[str, str]
    # Recovery metadata
    attempt_count: int
    failure_type: str
    recovery_decision: str
    retry_backoff_until: float
    correction_context: str
    step_traces: List[Dict[str, Any]]
    idempotency_key: str
    executed_side_effects: List[str]
    last_error: str
    last_error_signature: str
    error_signature_count: int

class WorkflowState(TypedDict):
    run_id: str
    user_id: str
    user_request: str
    run_title: str
    plan: Dict[str, Any]
    tasks: List[TaskDict]
    ready_queue: List[str]
    running_tasks: List[str]
    completed_tasks: List[str]
    failed_tasks: List[str]
    artifacts: List[str]
    notifications: List[str]
    debug_logs: List[str]
    final_result: Optional[str]
    status: str  # "pending", "running", "completed", "failed", "paused"


# ── Worker Templates ─────────────────────────────────────────────────

TEMPLATES_FILE = os.path.join(DATA_DIR, "agent_templates.json")

DEFAULT_WORKER_TEMPLATES = {
    "researcher": {
        "name": "Researcher Agent",
        "purpose": "Queries vector memory, crawls deep web search indices, and aggregates detailed reference facts.",
        "allowed_tools": ["web_search", "fetch_url", "memory_search"],
        "system_prompt_template": (
            "You are the Lead Researcher Agent working autonomously to gather exhaustive, verified information.\n"
            "You have access to tools — use them iteratively: search, fetch URLs, dig deeper into results, and aggregate evidence.\n"
            "Do NOT stop after a single search. Keep calling tools until you have thorough coverage of the topic.\n"
            "Task-specific context: {task_description}\n"
            "When fully satisfied with your research, produce your final structured output."
        ),
        "expected_output_schema": {
            "summary": "string",
            "sources": ["string"],
            "notes": ["string"]
        },
        "timeout": 300,
        "retry_limit": 3,
        "max_iterations": 15
    },
    "writer": {
        "name": "Writer Agent",
        "purpose": "Synthesizes raw factual notes into premium structured markdown reports.",
        "allowed_tools": ["artifact_read", "file_write"],
        "system_prompt_template": (
            "You are the Master Technical Writer producing exhaustive, structured markdown documents.\n"
            "Read any available artifacts first, then write the full content to disk using file_write.\n"
            "Do NOT produce a brief summary. Write a comprehensive, multi-section document.\n"
            "Task-specific context: {task_description}\n"
            "When the file is written, confirm with your final output."
        ),
        "expected_output_schema": {
            "title": "string",
            "markdown": "string"
        },
        "timeout": 180,
        "retry_limit": 2,
        "max_iterations": 10
    },
    "pdf_worker": {
        "name": "PDF Export Worker",
        "purpose": "Compiles reports and structured markdown into beautiful publication-grade PDF, Word, PowerPoint, and Excel documents.",
        "allowed_tools": ["artifact_read", "file_write", "pdf_export", "docx_export", "pptx_export", "excel_export"],
        "system_prompt_template": (
            "You are the Professional Document Layout Designer. Compile all drafts into publication-grade documents.\n"
            "Call the appropriate export tools (pdf_export, docx_export, etc.) with the full combined markdown content.\n"
            "Task-specific context: {task_description}\n"
            "Produce the requested file types, then report the generated file paths."
        ),
        "expected_output_schema": {
            "pdf_path": "string"
        },
        "timeout": 120,
        "retry_limit": 2,
        "max_iterations": 8
    },
    "email_worker": {
        "name": "Communications Worker",
        "purpose": "Dispatches summary notifications and attachments to user mailboxes or messaging platforms.",
        "allowed_tools": ["send_email"],
        "system_prompt_template": (
            "You are the Communications Officer. Dispatch summary notifications and attachments.\n"
            "Task-specific context: {task_description}\n"
            "Call the appropriate messaging tool and report the delivery status."
        ),
        "expected_output_schema": {
            "delivery_status": "string"
        },
        "timeout": 60,
        "retry_limit": 3,
        "max_iterations": 5
    },
    "browser": {
        "name": "Browser Automation Worker",
        "purpose": "Executes automated navigation, screenshot captures, and web crawls.",
        "allowed_tools": ["browser_automation"],
        "system_prompt_template": (
            "You are the Browser Navigator. Automate web pages step by step and report findings.\n"
            "Task-specific context: {task_description}"
        ),
        "expected_output_schema": {
            "page_details": "string"
        },
        "timeout": 300,
        "retry_limit": 2,
        "max_iterations": 20
    },
    "code_worker": {
        "name": "Execution Worker",
        "purpose": "Executes shell commands, scripts, docker, SSH operations, and file manipulations iteratively.",
        "allowed_tools": ["shell_exec", "file_write", "artifact_read"],
        "system_prompt_template": (
            "You are a Systems Execution Engineer running commands on the user's system autonomously.\n"
            "Use shell_exec to run commands one at a time. Observe the output before proceeding.\n"
            "Handle errors gracefully — if a command fails, diagnose and retry with a corrected command.\n"
            "You can run: bash commands, docker, docker-compose, pip, npm, ssh, curl, systemctl, etc.\n"
            "For long-running commands, use shell_exec with an appropriate timeout parameter.\n"
            "Task-specific context: {task_description}\n"
            "Work iteratively until the task is fully complete, then report results."
        ),
        "expected_output_schema": {
            "execution_output": "string",
            "summary": "string"
        },
        "timeout": 600,
        "retry_limit": 2,
        "max_iterations": 30
    }
}

def load_templates() -> dict:
    if not os.path.exists(TEMPLATES_FILE):
        _save(TEMPLATES_FILE, DEFAULT_WORKER_TEMPLATES)
        return DEFAULT_WORKER_TEMPLATES
    data = _load(TEMPLATES_FILE)
    if not isinstance(data, dict):
        return DEFAULT_WORKER_TEMPLATES
    return data

def save_templates(data: dict) -> None:
    _save(TEMPLATES_FILE, data)


# ── Tool Definitions ──────────────────────────────────────────────────

@tool
def web_search(query: str, max_results: int = 5) -> str:
    """Search the web for up-to-date facts on a specific query, specifying the number of results desired (max_results)."""
    prefs = load_prefs()
    provider = prefs.get("search_provider") or "tavily"
    
    if provider == "searxng":
        import httpx
        searxng_url = prefs.get("searxng_url") or "http://localhost:8080"
        searxng_url = searxng_url.rstrip("/")
        try:
            resp = httpx.get(
                f"{searxng_url}/",
                params={"q": query, "format": "json"},
                timeout=15.0
            )
            if resp.status_code != 200:
                resp = httpx.get(
                    f"{searxng_url}/search",
                    params={"q": query, "format": "json"},
                    timeout=10.0
                )
            if resp.status_code == 200:
                data = resp.json()
                raw_results = data.get("results", [])
                formatted_results = []
                for r in raw_results[:max_results]:
                    formatted_results.append({
                        "title": r.get("title") or "",
                        "url": r.get("url") or "",
                        "content": r.get("content") or r.get("snippet") or ""
                    })
                if formatted_results:
                    return json.dumps(formatted_results)
                return "Error: SearxNG search returned 0 results. The search engine might be rate-limited, experiencing connection issues, or blocking queries. Please proceed using your existing comprehensive knowledge, training, and logical deduction to draft the report contents."
            else:
                return f"Error: SearxNG search failed with status code {resp.status_code}. Please proceed using your existing comprehensive knowledge, training, and logical deduction to draft the report contents."
        except Exception as e:
            return f"Error: SearxNG search failed: {str(e)}. Please proceed using your existing comprehensive knowledge, training, and logical deduction to draft the report contents."

    # Tavily
    api_key = prefs.get("tavily_api_key") or os.getenv("TAVILY_API_KEY")
    if api_key:
        try:
            import os
            os.environ["TAVILY_API_KEY"] = api_key
            from langchain_community.tools.tavily_search import TavilySearchResults
            tavily = TavilySearchResults(max_results=max_results, tavily_api_key=api_key)
            results = tavily.invoke(query)
            if results:
                return json.dumps(results)
            return "Error: Tavily search returned 0 results. The API key might have run out of credits or hit its limit. Please proceed using your existing comprehensive knowledge, training, and logical deduction to draft the report contents."
        except Exception as e:
            return f"Error: Tavily search failed: {str(e)}. Please proceed using your existing comprehensive knowledge, training, and logical deduction to draft the report contents."
            
    return "Error: No search provider API keys are configured. Please proceed using your existing comprehensive knowledge, training, and logical deduction to draft the report contents."

@tool
def scrape_page(url: str) -> str:
    """Scrape a webpage and return clean text content, stripped of JavaScript, CSS, and HTML tags."""
    import httpx
    from html.parser import HTMLParser
    
    class TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text_parts = []
            self.in_ignored_tag = False
            self.ignored_tags = {"script", "style", "head", "meta", "link", "noscript", "svg"}

        def handle_starttag(self, tag, attrs):
            if tag in self.ignored_tags:
                self.in_ignored_tag = True

        def handle_endtag(self, tag):
            if tag in self.ignored_tags:
                self.in_ignored_tag = False

        def handle_data(self, data):
            if not self.in_ignored_tag:
                clean_data = data.strip()
                if clean_data:
                    self.text_parts.append(clean_data)

        def get_text(self):
            return "\n".join(self.text_parts)

    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }
        resp = httpx.get(url, headers=headers, follow_redirects=True, timeout=15.0)
        if resp.status_code == 200:
            parser = TextExtractor()
            parser.feed(resp.text)
            text = parser.get_text()
            if len(text) > 12000:
                text = text[:12000] + "\n\n[Content truncated due to length limits...]"
            return text if text.strip() else "Webpage returned no extractable text content."
        else:
            return f"Failed to retrieve page content. Status code: {resp.status_code}"
    except Exception as e:
        return f"Error occurred while scraping the page: {str(e)}"

@tool
def fetch_url(url: str) -> str:
    """Fetch the textual content of a specific webpage or URL, returning clean readable text rather than raw HTML. Automatically parses headers, paragraphs, lists, tables, and image links."""
    try:
        import httpx
        from html.parser import HTMLParser
        
        class SmartExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text_parts = []
                self.in_ignored = False
                self.ignored_tags = {"script", "style", "head", "meta", "link", "noscript", "svg", "nav", "footer"}
                
            def handle_starttag(self, tag, attrs):
                if tag in self.ignored_tags:
                    self.in_ignored = True
                elif not self.in_ignored:
                    if tag in ["h1", "h2", "h3", "h4", "h5", "h6"]:
                        self.text_parts.append("\n\n" + "#" * int(tag[1]) + " ")
                    elif tag in ["p", "div", "br", "tr", "li", "blockquote"]:
                        self.text_parts.append("\n")
                    elif tag == "img":
                        attr_dict = dict(attrs)
                        src = attr_dict.get("src", "")
                        alt = attr_dict.get("alt", "Image")
                        if src and not src.startswith("data:image"):
                            self.text_parts.append(f"\n[{alt}: {src}]\n")
                    elif tag == "td" or tag == "th":
                        self.text_parts.append(" | ")
                        
            def handle_endtag(self, tag):
                if tag in self.ignored_tags:
                    self.in_ignored = False
                elif not self.in_ignored:
                    if tag in ["h1", "h2", "h3", "h4", "h5", "h6", "p", "div", "tr", "li", "blockquote"]:
                        self.text_parts.append("\n")
                        
            def handle_data(self, data):
                if not self.in_ignored:
                    clean = data.strip()
                    if clean:
                        self.text_parts.append(clean + " ")
                        
            def get_text(self):
                import re
                raw = "".join(self.text_parts)
                # Collapse more than 2 newlines into 2
                return re.sub(r'\n{3,}', '\n\n', raw).strip()

        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }
        resp = httpx.get(url, headers=headers, follow_redirects=True, timeout=15.0)
        
        if resp.status_code != 200:
            return f"Fetch failed with HTTP {resp.status_code}"
            
        extractor = SmartExtractor()
        extractor.feed(resp.text)
        content = extractor.get_text()
        
        return content[:8000] if len(content) > 8000 else content

    except Exception as e:
        return f"Fetch failed: {str(e)}"

@tool
def escalate_to_supervisor(reason: str) -> str:
    """If you are completely blocked (e.g. missing a critical file that a previous agent was supposed to create), call this tool to instantly stop your work and escalate to the recovery supervisor. Provide a clear reason why the node needs to be restarted or investigated."""
    return f"ESCALATION_TRIGGERED: {reason}"

@tool
def memory_search(query: str) -> str:
    """Retrieve relevant memories matching the query from Qdrant vector storage."""
    try:
        results = search_memories("admin", query, limit=3)
        if not results:
            return "No vector memory entries found matching query."
        return json.dumps([r.payload for r in results])
    except Exception as e:
        return f"Memory search error: {str(e)}"

@tool
def artifact_read(filepath: str) -> str:
    """Read a text or markdown file from the active storage workspace or data directory."""
    try:
        sdir = active_storage_dir.get()
        if sdir:
            target = os.path.abspath(os.path.join(sdir, filepath.lstrip("/")))
            if target.startswith(os.path.abspath(sdir)) and os.path.exists(target):
                with open(target, "r") as f:
                    return f.read()
                    
        clean_path = os.path.basename(filepath)
        target = os.path.join(DATA_DIR, clean_path)
        if not os.path.exists(target):
            target = filepath
        with open(target, "r") as f:
            return f.read()
    except Exception as e:
        return f"Error reading file: {str(e)}"

@tool
def file_read(filepath: str) -> str:
    """Read the contents of any file on the system. Use absolute paths or relative paths. Useful for reading configuration files, code, or logs."""
    try:
        import os
        
        if not os.path.isabs(filepath):
            sdir = active_storage_dir.get() or os.getcwd()
            filepath = os.path.join(sdir, filepath)
            
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
            if len(content) > 15000:
                return content[:15000] + "\n\n...[FILE CONTENT TRUNCATED DUE TO LENGTH]..."
            return content
    except Exception as e:
        return f"Error reading file '{filepath}': {str(e)}"

@tool
def file_write(filepath: str, content: str) -> str:
    """Write text content to a specific file in the local data directory. Supports nested relative paths."""
    try:
        sdir = active_storage_dir.get()
        if not sdir:
            sdir = DATA_DIR
        target = os.path.abspath(os.path.join(sdir, filepath.lstrip("/")))
        if not target.startswith(os.path.abspath(sdir)):
            # fallback to base name if path traversal detected
            target = os.path.join(sdir, os.path.basename(filepath))
            
        os.makedirs(os.path.dirname(target), exist_ok=True)
        with open(target, "w") as f:
            f.write(content)
        return f"Successfully wrote file to: {target}"
    except Exception as e:
        return f"Error writing file: {str(e)}"

@tool
def directory_create(dirpath: str) -> str:
    """Create a new subdirectory inside the active workflow folder."""
    sdir = active_storage_dir.get()
    if not sdir:
        return "Error: No active workflow storage context."
    target = os.path.abspath(os.path.join(sdir, dirpath.lstrip("/")))
    if not target.startswith(os.path.abspath(sdir)):
        return "Error: Cannot create directory outside workflow storage."
    try:
        os.makedirs(target, exist_ok=True)
        return f"Successfully created directory at: {dirpath}"
    except Exception as e:
        return f"Failed to create directory: {str(e)}"

@tool
def pdf_export(markdown_content: str, filename: str) -> str:
    """Render structured markdown text into a publication-grade PDF file using ReportLab."""
    try:
        def clean_unicode_text(text: str) -> str:
            replacements = {
                "\u2010": "-",  # Hyphen
                "\u2011": "-",  # Non-breaking hyphen
                "\u2012": "-",  # Figure dash
                "\u2013": "-",  # En dash (–)
                "\u2014": "--", # Em dash (—)
                "\u2015": "--", # Horizontal bar
                "\u2212": "-",  # Minus sign (−)
                "\xa0": " ",    # Non-breaking space
                "\u2018": "'",  # Left single quote
                "\u2019": "'",  # Right single quote
                "\u201c": '"',  # Left double quote
                "\u201d": '"',  # Right double quote
                "\u2022": "*",  # Bullet
                "■": "-",       # Fallback replacement
            }
            for k, v in replacements.items():
                text = text.replace(k, v)
            return text

        markdown_content = clean_unicode_text(markdown_content)

        clean_name = os.path.basename(filename)
        if not clean_name.endswith(".pdf"):
            clean_name += ".pdf"

        # Dynamically rename generic filenames to match active workflow title/topic
        generic_names = ["generated", "compiled", "report", "final", "document", "output", "booklet", "pdf", "export"]
        name_no_ext = os.path.splitext(clean_name)[0].lower()
        if name_no_ext in generic_names or not name_no_ext.strip():
            wf_title = active_workflow_title.get()
            if wf_title:
                clean_name = _re.sub(r'[^a-zA-Z0-9_\-]', '_', wf_title).strip("_")[:50] + ".pdf"

        uploads_dir = active_storage_dir.get() or os.path.join(DATA_DIR, "uploads")
        os.makedirs(uploads_dir, exist_ok=True)
        pdf_path = os.path.join(uploads_dir, clean_name)

        doc = SimpleDocTemplate(pdf_path, pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
        styles = getSampleStyleSheet()

        # --- Custom styles ---
        s_title = ParagraphStyle('S_Title', parent=styles['Heading1'], fontSize=22, leading=28,
                                  textColor=colors.HexColor('#272757'), spaceAfter=14, keepWithNext=True)
        s_h2 = ParagraphStyle('S_H2', parent=styles['Heading2'], fontSize=16, leading=20,
                               textColor=colors.HexColor('#505081'), spaceBefore=14, spaceAfter=8, keepWithNext=True)
        s_h3 = ParagraphStyle('S_H3', parent=styles['Heading3'], fontSize=13, leading=17,
                               textColor=colors.HexColor('#505081'), spaceBefore=10, spaceAfter=6, keepWithNext=True)
        s_h4 = ParagraphStyle('S_H4', parent=styles['Heading4'], fontSize=11, leading=15,
                               textColor=colors.HexColor('#8686AC'), spaceBefore=8, spaceAfter=4, keepWithNext=True)
        s_body = ParagraphStyle('S_Body', parent=styles['Normal'], fontSize=10, leading=15,
                                 textColor=colors.HexColor('#1F2937'), spaceAfter=6)
        s_bullet = ParagraphStyle('S_Bullet', parent=s_body, leftIndent=18, bulletIndent=6)
        s_quote = ParagraphStyle('S_Quote', parent=s_body, leftIndent=20, textColor=colors.HexColor('#6B7280'),
                                  borderPadding=4, fontName='Helvetica-Oblique')
        s_tbl_header = ParagraphStyle('S_TblH', parent=styles['Normal'], fontSize=9, leading=12,
                                       textColor=colors.HexColor('#FFFFFF'), fontName='Helvetica-Bold')
        s_tbl_cell = ParagraphStyle('S_TblC', parent=styles['Normal'], fontSize=9, leading=12,
                                     textColor=colors.HexColor('#374151'))

        def _md_inline(t):
            """Convert inline markdown (bold, italic, links, code) to ReportLab markup robustly."""
            # 1. Escape XML characters first
            t = t.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            t = t.replace('&amp;amp;', '&amp;').replace('&amp;lt;', '&lt;').replace('&amp;gt;', '&gt;')
            
            # 2. Extract code blocks to prevent formatting inside them
            code_blocks = []
            def _sub_code(m):
                code_blocks.append(m.group(1))
                return f"XYZCODEBLOCK{len(code_blocks)-1}XYZ"
                
            t = _re.sub(r'`(.+?)`', _sub_code, t)
            
            # 3. Process links, bold+italic combos, then bold, then italics
            t = _re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'<font color="#505081"><u>\1</u></font>', t)
            # Bold+italic combo (***text*** or ___text___) — must come BEFORE bold/italic
            t = _re.sub(r'\*\*\*(.+?)\*\*\*', r'<b><i>\1</i></b>', t)
            t = _re.sub(r'___(.+?)___', r'<b><i>\1</i></b>', t)
            t = _re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', t)
            t = _re.sub(r'__(.+?)__', r'<b>\1</b>', t)
            t = _re.sub(r'\*(.+?)\*', r'<i>\1</i>', t)
            t = _re.sub(r'_(.+?)_', r'<i>\1</i>', t)
            
            # 4. Re-insert code blocks formatted for ReportLab
            for idx, cv in enumerate(code_blocks):
                t = t.replace(f"XYZCODEBLOCK{idx}XYZ", f'<font face="Courier" size="8" color="#505081">{cv}</font>')
                
            t = t.replace('■', '-')
            
            # 5. Repair any improperly nested XML tags (e.g. <b><i>...</b></i> → <b><i>...</i></b>)
            t = _repair_xml_nesting(t)
            return t

        def _repair_xml_nesting(text):
            """Fix incorrectly nested XML/HTML tags that crash ReportLab's strict parser."""
            import re as _rr
            tag_pattern = _rr.compile(r'<(/?)(\w+)(?:\s[^>]*)?>') 
            tokens = tag_pattern.split(text)
            # Rebuild with a stack-based approach
            stack = []
            result_parts = []
            pos = 0
            for m in tag_pattern.finditer(text):
                # Add text before this tag
                result_parts.append(text[pos:m.start()])
                pos = m.end()
                is_close = m.group(1) == '/'
                tag_name = m.group(2).lower()
                
                # Skip self-closing or non-formatting tags
                if tag_name not in ('b', 'i', 'u', 'font'):
                    result_parts.append(m.group(0))
                    continue
                    
                if not is_close:
                    stack.append(tag_name)
                    result_parts.append(m.group(0))
                else:
                    if tag_name in stack:
                        # Close tags in reverse order up to the matching one
                        to_reopen = []
                        while stack and stack[-1] != tag_name:
                            popped = stack.pop()
                            result_parts.append(f'</{popped}>')
                            to_reopen.append(popped)
                        if stack:
                            stack.pop()
                        result_parts.append(f'</{tag_name}>')
                        # Re-open any tags we had to close early
                        for reopened in reversed(to_reopen):
                            stack.append(reopened)
                            result_parts.append(f'<{reopened}>')
                    else:
                        pass  # Ignore orphaned closing tag
            # Add remaining text
            result_parts.append(text[pos:])
            return ''.join(result_parts)

        def _safe_paragraph(text, style):
            """Wrap Paragraph creation with error handling so one bad line doesn't crash the whole PDF."""
            try:
                return Paragraph(text, style)
            except Exception:
                # Strip ALL markup and retry as plain text
                clean = _re.sub(r'<[^>]+>', '', text)
                try:
                    return Paragraph(clean, style)
                except Exception:
                    return Paragraph("(Content could not be rendered)", style)

        story = []
        story.append(Paragraph("KOKOMI STRATEGIST OS — INTELLIGENCE REPORT", s_title))
        story.append(Spacer(1, 8))

        lines = markdown_content.split("\n")
        i = 0
        while i < len(lines):
            raw = lines[i]
            stripped = raw.strip()

            # Empty line
            if not stripped:
                i += 1
                continue

            # Horizontal rule
            if stripped in ('---', '***', '___'):
                story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor('#D1D5DB'),
                                         spaceBefore=8, spaceAfter=8))
                i += 1
                continue

            # Table detection (| ... | ... |)
            if '|' in stripped and stripped.startswith('|'):
                table_lines = []
                while i < len(lines) and lines[i].strip().startswith('|'):
                    table_lines.append(lines[i].strip())
                    i += 1
                # Parse table
                rows = []
                for tl in table_lines:
                    cells = [c.strip() for c in tl.strip('|').split('|')]
                    if all(set(c.strip()) <= set('-| :') for c in cells):
                        continue  # Skip separator row
                    rows.append(cells)
                if rows:
                    col_count = max(len(r) for r in rows)
                    tbl_data = []
                    for ri, row in enumerate(rows):
                        while len(row) < col_count:
                            row.append('')
                        style = s_tbl_header if ri == 0 else s_tbl_cell
                        tbl_data.append([_safe_paragraph(_md_inline(c), style) for c in row])
                    col_w = (letter[0] - 80) / col_count
                    t = Table(tbl_data, colWidths=[col_w] * col_count, repeatRows=1)
                    t.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#505081')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#D1D5DB')),
                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor('#F9FAFB'), colors.white]),
                        ('TOPPADDING', (0, 0), (-1, -1), 4),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                        ('LEFTPADDING', (0, 0), (-1, -1), 6),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                    ]))
                    story.append(Spacer(1, 6))
                    story.append(t)
                    story.append(Spacer(1, 6))
            # Page Break
            if stripped in ('<!-- pagebreak -->', '\\pagebreak', r'\pagebreak', '[pagebreak]'):
                story.append(PageBreak())
                i += 1
                continue
                
            # Full-line Image Detection (![Alt](URL))
            img_match = _re.match(r'^!\[(.*?)\]\((.*?)\)$', stripped)
            if img_match:
                img_url = img_match.group(2).strip()
                try:
                    import requests
                    import tempfile
                    import uuid
                    
                    local_img = img_url
                    if img_url.startswith("http"):
                        # Download it dynamically!
                        img_resp = requests.get(img_url, timeout=15)
                        if img_resp.status_code == 200:
                            tmp_path = os.path.join(tempfile.gettempdir(), f"img_{uuid.uuid4().hex}.png")
                            with open(tmp_path, "wb") as f:
                                f.write(img_resp.content)
                            local_img = tmp_path
                            
                    if os.path.exists(local_img):
                        from reportlab.platypus import Image as RLImage
                        img = RLImage(local_img)
                        max_w = letter[0] - 80
                        if img.drawWidth > max_w:
                            ratio = max_w / float(img.drawWidth)
                            img.drawWidth = max_w
                            img.drawHeight = img.drawHeight * ratio
                        story.append(Spacer(1, 10))
                        story.append(img)
                        story.append(Spacer(1, 10))
                        i += 1
                        continue
                except Exception as e:
                    print(f"Failed to embed image {img_url}: {e}")
                    # If it fails, fallback to rendering the raw alt-text below

            # Multi-line code block detection
            if stripped.startswith('```'):
                code_lang = stripped[3:].strip().lower()
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith('```'):
                    code_lines.append(lines[i]) # Keep raw spacing!
                    i += 1
                if i < len(lines):
                    i += 1 # Skip closing ```
                
                code_text = "\n".join(code_lines)
                
                # --- MERMAID INTERCEPTOR ---
                if code_lang == 'mermaid':
                    try:
                        import urllib.parse
                        import requests
                        import tempfile
                        import uuid
                        
                        encoded = urllib.parse.quote(code_text)
                        url = f"https://quickchart.io/mermaid?graph={encoded}"
                        img_resp = requests.get(url, timeout=15)
                        
                        if img_resp.status_code == 200:
                            tmp_path = os.path.join(tempfile.gettempdir(), f"mermaid_{uuid.uuid4().hex}.png")
                            with open(tmp_path, "wb") as f:
                                f.write(img_resp.content)
                            from reportlab.platypus import Image as RLImage
                            img = RLImage(tmp_path)
                            max_w = letter[0] - 80
                            if img.drawWidth > max_w:
                                ratio = max_w / float(img.drawWidth)
                                img.drawWidth = max_w
                                img.drawHeight = img.drawHeight * ratio
                            story.append(Spacer(1, 10))
                            story.append(img)
                            story.append(Spacer(1, 10))
                            continue
                    except Exception as e:
                        print(f"Mermaid rendering failed: {e}")
                
                # Format code content
                # Escape XML characters for ReportLab Paragraph
                code_text = code_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                
                s_code_block = ParagraphStyle('S_CodeBlock', parent=styles['Normal'],
                                               fontName='Courier', fontSize=8, leading=11,
                                               textColor=colors.HexColor('#0F0E47'))
                
                # Wrap inside a Table to give it a gorgeous light-grey background and nice borders!
                p_code = Preformatted(code_text, s_code_block)
                t_code = Table([[p_code]], colWidths=[letter[0] - 80])
                t_code.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#F4F4F7')),
                    ('BOX', (0, 0), (-1, -1), 0.5, colors.HexColor('#8686AC')),
                    ('TOPPADDING', (0, 0), (-1, -1), 6),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                    ('LEFTPADDING', (0, 0), (-1, -1), 8),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                ]))
                story.append(Spacer(1, 4))
                story.append(t_code)
                story.append(Spacer(1, 4))
                continue

            # Headings
            if stripped.startswith('#### '):
                story.append(_safe_paragraph(_md_inline(stripped[5:]), s_h4))
            elif stripped.startswith('### '):
                story.append(_safe_paragraph(_md_inline(stripped[4:]), s_h3))
            elif stripped.startswith('## '):
                story.append(_safe_paragraph(_md_inline(stripped[3:]), s_h2))
            elif stripped.startswith('# '):
                story.append(_safe_paragraph(_md_inline(stripped[2:]), s_title))
            # Blockquote
            elif stripped.startswith('> '):
                story.append(_safe_paragraph(_md_inline(stripped[2:]), s_quote))
            # Bullets
            elif stripped.startswith('- ') or stripped.startswith('* '):
                story.append(_safe_paragraph("• " + _md_inline(stripped[2:]), s_bullet))
            # Numbered list
            elif _re.match(r'^\d+\.\s', stripped):
                m = _re.match(r'^(\d+\.)\s(.*)', stripped)
                story.append(_safe_paragraph(f"{m.group(1)} {_md_inline(m.group(2))}", s_bullet))
            # Regular paragraph
            else:
                story.append(_safe_paragraph(_md_inline(stripped), s_body))

            i += 1

        doc.build(story)
        return f"Successfully rendered PDF at: {pdf_path}"
    except Exception as e:
        return f"PDF generation failed: {str(e)}"


@tool
def docx_export(markdown_content: str, filename: str) -> str:
    """Render structured markdown text into a beautifully styled Microsoft Word DOCX document."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        clean_name = os.path.basename(filename)
        if not clean_name.endswith(".docx"):
            clean_name += ".docx"
            
        # Naming override if generic
        generic_names = ["generated", "compiled", "report", "final", "document", "output", "booklet", "docx", "export"]
        name_no_ext = os.path.splitext(clean_name)[0].lower()
        if name_no_ext in generic_names or not name_no_ext.strip():
            wf_title = active_workflow_title.get()
            if wf_title:
                clean_name = _re.sub(r'[^a-zA-Z0-9_\-]', '_', wf_title).strip("_")[:50] + ".docx"

        uploads_dir = active_storage_dir.get() or os.path.join(DATA_DIR, "uploads")
        os.makedirs(uploads_dir, exist_ok=True)
        docx_path = os.path.join(uploads_dir, clean_name)
        
        doc = Document()
        
        # Apply standard styles / apple HIG colors
        c_midnight = RGBColor(0x27, 0x27, 0x57)
        c_indigo = RGBColor(0x50, 0x50, 0x81)
        c_lavender = RGBColor(0x86, 0x86, 0xAC)
        
        def add_formatted_text(paragraph, text, size=10.5, color=None, force_bold=False):
            import re as _re_clean
            parts = _re_clean.split(r'(\*\*.*?\*\*)', text)
            for part in parts:
                if not part:
                    continue
                if part.startswith('**') and part.endswith('**'):
                    clean_part = part[2:-2].replace('*', '')
                    run = paragraph.add_run(clean_part)
                    run.bold = True
                else:
                    clean_part = part.replace('*', '')
                    run = paragraph.add_run(clean_part)
                    if force_bold:
                        run.bold = True
                run.font.name = 'Arial'
                run.font.size = Pt(size)
                if color:
                    run.font.color.rgb = color

        lines = markdown_content.split("\n")
        in_code_block = False
        in_bullet_list = False
        
        for line in lines:
            stripped = line.strip()
            if stripped.startswith("```"):
                in_code_block = not in_code_block
                continue
                
            if in_code_block:
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Inches(0.5)
                run = p.add_run(line)
                run.font.name = 'Consolas'
                run.font.size = Pt(9.5)
                run.font.color.rgb = RGBColor(100, 100, 100)
                continue
                
            # Headers
            if stripped.startswith("# "):
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(12)
                p.paragraph_format.space_after = Pt(6)
                add_formatted_text(p, stripped[2:], size=18, color=c_midnight, force_bold=True)
                in_bullet_list = False
            elif stripped.startswith("## "):
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(10)
                p.paragraph_format.space_after = Pt(4)
                add_formatted_text(p, stripped[3:], size=14, color=c_indigo, force_bold=True)
                in_bullet_list = False
            elif stripped.startswith("### "):
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(8)
                p.paragraph_format.space_after = Pt(4)
                add_formatted_text(p, stripped[4:], size=12, color=c_lavender, force_bold=True)
                in_bullet_list = False
            # Bullet items
            elif stripped.startswith("- ") or stripped.startswith("* "):
                p = doc.add_paragraph(style='List Bullet')
                p.paragraph_format.space_after = Pt(2)
                add_formatted_text(p, stripped[2:])
                in_bullet_list = True
            elif stripped.startswith("1. ") or _re.match(r'^\d+\.\s', stripped):
                match = _re.match(r'^(\d+\.\s)', stripped)
                prefix_len = len(match.group(1))
                p = doc.add_paragraph(style='List Number')
                p.paragraph_format.space_after = Pt(2)
                add_formatted_text(p, stripped[prefix_len:])
                in_bullet_list = True
            elif not stripped:
                in_bullet_list = False
                continue
            else:
                p = doc.add_paragraph()
                p.paragraph_format.space_after = Pt(4)
                add_formatted_text(p, line)
                in_bullet_list = False
                
        doc.save(docx_path)
        return f"Successfully compiled and saved beautiful Word Document (.docx) at: {docx_path}"
    except Exception as e:
        return f"Failed to generate Word document: {str(e)}"


@tool
def pptx_export(markdown_content: str, filename: str) -> str:
    """Render structured markdown text (separated by '---' or Heading 2s) into a premium PowerPoint Presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        clean_name = os.path.basename(filename)
        if not clean_name.endswith(".pptx"):
            clean_name += ".pptx"
            
        # Naming override if generic
        generic_names = ["generated", "compiled", "report", "final", "document", "output", "booklet", "pptx", "presentation", "export"]
        name_no_ext = os.path.splitext(clean_name)[0].lower()
        if name_no_ext in generic_names or not name_no_ext.strip():
            wf_title = active_workflow_title.get()
            if wf_title:
                clean_name = _re.sub(r'[^a-zA-Z0-9_\-]', '_', wf_title).strip("_")[:50] + ".pptx"

        uploads_dir = active_storage_dir.get() or os.path.join(DATA_DIR, "uploads")
        os.makedirs(uploads_dir, exist_ok=True)
        pptx_path = os.path.join(uploads_dir, clean_name)
        
        prs = Presentation()
        
        c_midnight = RGBColor(0x27, 0x27, 0x57)
        c_indigo = RGBColor(0x50, 0x50, 0x81)
        c_lavender = RGBColor(0x86, 0x86, 0xAC)
        c_dark = RGBColor(0x11, 0x18, 0x27)
        
        import re as _re_clean
        def clean_title(text: str) -> str:
            # Strip slide titles of any bold/italic decorators
            text = _re_clean.sub(r'\*\*(.*?)\*\*', r'\1', text)
            text = _re_clean.sub(r'\*(.*?)\*', r'\1', text)
            return text.strip()
            
        slides_text = []
        if "---" in markdown_content:
            slides_text = [s.strip() for s in markdown_content.split("---")]
        else:
            current = []
            for line in markdown_content.split("\n"):
                if line.strip().startswith("## ") and current:
                    slides_text.append("\n".join(current))
                    current = [line]
                else:
                    current.append(line)
            if current:
                slides_text.append("\n".join(current))
                
        # Title Slide (First slide)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        wf_title = active_workflow_title.get() or "Professional Presentation"
        title.text = wf_title
        title.text_frame.paragraphs[0].font.color.rgb = c_midnight
        title.text_frame.paragraphs[0].font.name = 'Arial'
        subtitle.text = "Generated Automatically by Antigravity OS Compiler"
        subtitle.text_frame.paragraphs[0].font.color.rgb = c_indigo
        
        # Body Slides
        bullet_slide_layout = prs.slide_layouts[1]
        
        for stext in slides_text:
            lines = [l.strip() for l in stext.split("\n") if l.strip()]
            if not lines:
                continue
                
            slide_title = "Overview"
            bullet_points = []
            
            for line in lines:
                if line.startswith("# ") or line.startswith("## ") or line.startswith("### "):
                    slide_title = clean_title(line.lstrip("#").strip())
                elif line.startswith("- ") or line.startswith("* "):
                    bullet_points.append(line[2:])
                else:
                    if len(line) > 5:
                        bullet_points.append(line)
                        
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Slide Title
            tx_title = slide.shapes.title
            tx_title.text = slide_title
            tx_title.text_frame.paragraphs[0].font.color.rgb = c_midnight
            tx_title.text_frame.paragraphs[0].font.name = 'Arial'
            tx_title.text_frame.paragraphs[0].font.bold = True
            
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            
            for idx, pt in enumerate(bullet_points[:6]):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.level = 0
                
                parts = _re_clean.split(r'(\*\*.*?\*\*)', pt)
                for part in parts:
                    if not part:
                        continue
                    if part.startswith('**') and part.endswith('**'):
                        clean_part = part[2:-2].replace('*', '')
                        run = p.add_run()
                        run.text = clean_part
                        run.font.bold = True
                    else:
                        clean_part = part.replace('*', '')
                        run = p.add_run()
                        run.text = clean_part
                        
                    run.font.size = Pt(15)
                    run.font.name = 'Arial'
                    run.font.color.rgb = c_dark
                    
        prs.save(pptx_path)
        return f"Successfully compiled and saved beautiful PowerPoint (.pptx) at: {pptx_path}"
    except Exception as e:
        return f"Failed to generate PowerPoint: {str(e)}"


@tool
def excel_export(table_data_json: str, filename: str) -> str:
    """Render structured tabular JSON or CSV data into a beautiful, styled Excel spreadsheet (.xlsx)."""
    try:
        import json
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
        
        clean_name = os.path.basename(filename)
        if not clean_name.endswith(".xlsx"):
            clean_name += ".xlsx"
            
        # Naming override if generic
        generic_names = ["generated", "compiled", "report", "final", "document", "output", "booklet", "excel", "xlsx", "export"]
        name_no_ext = os.path.splitext(clean_name)[0].lower()
        if name_no_ext in generic_names or not name_no_ext.strip():
            wf_title = active_workflow_title.get()
            if wf_title:
                clean_name = _re.sub(r'[^a-zA-Z0-9_\-]', '_', wf_title).strip("_")[:50] + ".xlsx"

        uploads_dir = active_storage_dir.get() or os.path.join(DATA_DIR, "uploads")
        os.makedirs(uploads_dir, exist_ok=True)
        excel_path = os.path.join(uploads_dir, clean_name)
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Report Data"
        
        data = []
        try:
            data = json.loads(table_data_json)
        except Exception:
            lines = [l.strip() for l in table_data_json.split("\n") if l.strip()]
            for line in lines:
                data.append([c.strip() for c in line.split(",")])
                
        if not data:
            return "Failed to generate Excel: Provided data was empty or invalid."
            
        rows = []
        headers = []
        if isinstance(data, list) and len(data) > 0:
            if isinstance(data[0], dict):
                headers = list(data[0].keys())
                rows.append(headers)
                for obj in data:
                    rows.append([obj.get(h, "") for h in headers])
            else:
                rows = data
                if len(rows) > 0:
                    headers = rows[0]
        else:
            return "Failed to generate Excel: Unrecognized data format."
            
        font_header = Font(name="Arial", size=11, bold=True, color="FFFFFF")
        fill_header = PatternFill(start_color="272757", end_color="272757", fill_type="solid")
        
        font_cell = Font(name="Arial", size=10)
        fill_zebra = PatternFill(start_color="F9FAFB", end_color="F9FAFB", fill_type="solid")
        
        thin_border = Border(
            left=Side(style='thin', color='E5E7EB'),
            right=Side(style='thin', color='E5E7EB'),
            top=Side(style='thin', color='E5E7EB'),
            bottom=Side(style='thin', color='E5E7EB')
        )
        
        for r_idx, row in enumerate(rows, 1):
            ws.append(row)
            is_header = (r_idx == 1)
            for c_idx in range(1, len(row) + 1):
                cell = ws.cell(row=r_idx, column=c_idx)
                cell.border = thin_border
                if is_header:
                    cell.font = font_header
                    cell.fill = fill_header
                    cell.alignment = Alignment(horizontal="center", vertical="center")
                else:
                    cell.font = font_cell
                    if r_idx % 2 == 0:
                        cell.fill = fill_zebra
                        
        for col in ws.columns:
            max_len = 0
            col_letter = get_column_letter(col[0].column)
            for cell in col:
                val_str = str(cell.value or '')
                if len(val_str) > max_len:
                    max_len = len(val_str)
            ws.column_dimensions[col_letter].width = max(max_len + 3, 12)
            
        wb.save(excel_path)
        return f"Successfully compiled and saved beautiful Excel (.xlsx) at: {excel_path}"
    except Exception as e:
        return f"Failed to generate Excel: {str(e)}"


@tool
def send_email(to_email: str, subject: str, body: str, attachment_path: Optional[str] = None) -> str:
    """Send out an email summary report to the user mailbox (includes attachments)."""
    try:
        # High fidelity simulate
        log_line = f"[{datetime.datetime.now().isoformat()}] Sending email to: {to_email}\nSubject: {subject}\nBody:\n{body}\nAttachment: {attachment_path}\n"
        logs_file = os.path.join(DATA_DIR, "email_logs.txt")
        with open(logs_file, "a") as f:
            f.write(log_line)
        return f"Email successfully dispatched to {to_email}. Delivery Status: SENT. Logged in data/email_logs.txt."
    except Exception as e:
        return f"Email transmission failed: {str(e)}"

@tool
def browser_automation(action: str, target: str) -> str:
    """Execute browser navigation or action simulation."""
    return f"Simulated browser automation task '{action}' on {target} completed successfully."

@tool
def shell_exec(command: str, timeout: int = 60) -> str:
    """Execute a shell command on the local system. Returns stdout+stderr combined. Use for running scripts, docker, ssh, pip installs, system operations. timeout is max seconds to wait (default 60, max 300)."""
    import subprocess
    from app.workflow import active_storage_dir
    timeout = min(max(timeout, 5), 300)
    try:
        sdir = active_storage_dir.get()
        if not sdir or not os.path.exists(sdir):
            sdir = None
            
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=timeout,
            cwd=sdir
        )
        output = ""
        if result.stdout:
            output += result.stdout
        if result.stderr:
            output += ("\n[stderr]\n" + result.stderr) if result.stdout else result.stderr
        if not output.strip():
            output = f"Command exited with code {result.returncode} (no output)."
        return output[:8000]  # cap to prevent context overflow
    except subprocess.TimeoutExpired:
        return f"Command timed out after {timeout}s."
    except Exception as e:
        return f"shell_exec error: {e}"


# ── Multi-Agent Routing Classifier ───────────────────────────────────

# Lightweight keyword-based pre-filter (cheap, catches obvious cases)
_WORKFLOW_KEYWORDS = [
    # Research & reporting
    "research", "write a report", "compile", "pdf", "generate report",
    "write a guide", "write a book", "write a manual", "write a whitepaper",
    # Execution & automation
    "install", "deploy", "docker", "docker-compose", "ssh", "run", "execute",
    "configure", "set up", "setup", "create a server", "start a service",
    "restart", "update packages", "pip install", "npm install",
    # Notifications
    "email me", "send email", "notify me", "whatsapp", "slack", "message me",
    "tell me when", "let me know when",
    # Multi-step signals
    "workflow", "multi-agent", "step by step", "and then",
    "save to file", "atlas", "automate",
]

def is_complex_workflow(request: str) -> bool:
    """Route to workflow engine vs simple chat. Uses keyword pre-filter for speed."""
    req_lower = request.lower()
    # Quick keyword pass
    if any(kw in req_lower for kw in _WORKFLOW_KEYWORDS):
        return True
    # If request is long and imperative (sounds like instructions), treat as workflow
    words = req_lower.split()
    if len(words) >= 10 and any(req_lower.startswith(v) for v in (
        "can you ", "could you ", "please ", "i want you to ", "i need you to ",
        "go ahead and ", "make ", "create ", "build ", "find ", "get me "
    )):
        return True
    return False


# ── Supervisor Planner & Plan Validation ─────────────────────────────

class TaskPlanSchema(BaseModel):
    run_title: str = Field(description="Exhaustive high-level title of the execution run.")
    goal: str = Field(description="High level outcome objective of the entire run.")
    tasks: List[Dict[str, Any]] = Field(description="The sequential/parallel tasks list to be executed.")

def validate_plan(plan: dict) -> bool:
    """Ensures structured task plan maps expected key formats."""
    if not isinstance(plan, dict):
        return False
    if "run_title" not in plan or "tasks" not in plan:
        return False
    for task in plan["tasks"]:
        required_keys = ["task_id", "title", "worker_type", "depends_on", "allowed_tools"]
        if not all(k in task for k in required_keys):
            return False
    return True

async def run_supervisor_planner(user_request: str) -> dict:
    """Invokes supervisor model to create high-fidelity task dependencies plan."""
    prefs = load_prefs()
    llm = get_atlas_llm(prefs, streaming=False)
    
    templates = load_templates()
    workers_desc = ""
    for idx, (worker_key, tpl) in enumerate(templates.items(), 1):
        name = tpl.get("name") or worker_key
        purpose = tpl.get("purpose") or ""
        allowed_tools = tpl.get("allowed_tools") or []
        workers_desc += f"{idx}. {worker_key} (allowed_tools: {allowed_tools}) - {purpose}\n"

    # Helper to find best matching worker type
    def get_best_worker(preferred: str, fallback_tool: str = None) -> str:
        if preferred in templates:
            return preferred
        if fallback_tool:
            for k, tpl in templates.items():
                if fallback_tool in tpl.get("allowed_tools", []):
                    return k
        return list(templates.keys())[0] if templates else preferred

    best_research_worker = get_best_worker("researcher", "web_search")
    best_writer_worker = get_best_worker("writer", "file_write")
    best_pdf_worker = get_best_worker("pdf_worker", "pdf_export")
    best_email_worker = get_best_worker("email_worker", "send_email")
    best_code_worker = get_best_worker("code_worker", "shell_exec")

    prompt = (
        "You are the Top-Level Workflow Supervisor. Autonomously decompose the user's objective into "
        "well-structured sequential and parallel tasks for specialized worker agents.\n"
        "Workers available:\n"
        f"{workers_desc}\n"
        "--- GENERAL PURPOSE AGENTIC SYSTEM ---\n"
        "This system is a general-purpose autonomous agent. It can run ANYTHING: research, writing, coding, "
        "DevOps, system administration, docker, SSH, notifications, file operations, and more.\n"
        "Each worker runs a ReAct loop (think → act → observe → repeat) and can call tools multiple times until done.\n\n"
        "--- DEVOPS / AUTOMATION / SHELL TASKS ---\n"
        f"For tasks involving commands, scripts, docker, SSH, installations, system config: use '{best_code_worker}' with 'shell_exec'.\n"
        "If an MCP tool is available (e.g. a shell MCP for a remote server, or a WhatsApp MCP): include it in 'allowed_tools'.\n"
        "Break complex operations into logical phases: e.g. 'Deploy containers', 'Install dependencies', 'Configure service', 'Notify user'.\n\n"
        "--- MASSIVE REPORT STRATEGY (40-50+ PAGES) ---\n"
        "If the user request asks for a very long, comprehensive, or highly detailed report/manual, "
        "break the writing phase into MULTIPLE sequential tasks (e.g., 'Write Chapter 1', 'Write Chapter 2', ...).\n"
        f"- Each writing task writes its output to a file (ch1.md, ch2.md, etc.) using 'file_write'.\n"
        f"- Plan a final '{best_pdf_worker}' task to compile all files into a PDF.\n\n"
        "--- NOTIFICATIONS ---\n"
        f"CRITICAL: ONLY add a final notifications / '{best_email_worker}' task if the user explicitly requested email or notifications in their request. DO NOT add any email, WhatsApp, or Slack notifications if the user did not explicitly request them.\n\n"
        f"User Request: {user_request}\n\n"
        "IMPORTANT: In 'allowed_tools', include any MCP tool names that are relevant (e.g. remote shell tools, messaging tools). "
        "The system will automatically route them to connected MCP servers.\n\n"
        "Respond ONLY with a valid JSON object. No markdown, no explanation.\n"
        "The worker_type and allowed_tools MUST match what is actually needed for the task — use researcher/writer/pdf_worker for research/reports, code_worker/shell_exec for system tasks, etc.\n"
        "JSON SCHEMA (this is just a format example, adapt worker_type and tools to fit the actual request):\n"
        "{\n"
        '  "run_title": "Concise title of the plan",\n'
        '  "goal": "High-level outcome objective",\n'
        '  "tasks": [\n'
        '    {\n'
        '      "task_id": "t1",\n'
        '      "title": "Task title",\n'
        '      "description": "Detailed instructions for this specific worker — be specific about what to do and what success looks like",\n'
        '      "worker_type": "<pick the right worker for this task>",\n'
        '      "depends_on": [],\n'
        '      "allowed_tools": ["<tools this worker needs>"],\n'
        '      "success_criteria": "Clear definition of done"\n'
        "    }\n"
        "  ]\n"
        "}"
    )
    
    try:
        response = await llm.ainvoke([HumanMessage(content=prompt)])
        raw_text = response.content.strip()
        # Parse JSON
        if "```json" in raw_text:
            raw_text = raw_text.split("```json")[1].split("```")[0].strip()
        elif "```" in raw_text:
            raw_text = raw_text.split("```")[1].split("```")[0].strip()
            
        # Robust JSON repair
        import re
        def repair_json_string(s: str) -> str:
            s = s.strip()
            # Find first brace and last brace
            first_brace = s.find('{')
            last_brace = s.rfind('}')
            if first_brace != -1 and last_brace != -1:
                s = s[first_brace:last_brace+1]
            # Remove single-line comments
            s = re.sub(r'^\s*//.*$', '', s, flags=re.MULTILINE)
            s = re.sub(r'\s*//.*$', '', s)
            # Remove block comments
            s = re.sub(r'/\*.*?\*/', '', s, flags=re.DOTALL)
            # Remove trailing commas inside lists/objects
            s = re.sub(r',\s*([\]}])', r'\1', s)
            return s.strip()

        cleaned_text = repair_json_string(raw_text)
        try:
            parsed_plan = json.loads(cleaned_text)
        except Exception as e_json:
            print(f"Standard JSON parse failed, trying advanced recovery: {e_json}")
            try:
                # Level 2 recovery: normalize key quotes and delimiters
                temp = cleaned_text
                temp = re.sub(r"(?<=[\{\[,:])\s*'", '"', temp)
                temp = re.sub(r"'\s*(?=[\}\],:])", '"', temp)
                temp = re.sub(r"(?<=[\{\,])\s*([a-zA-Z0-9_]+)\s*:", r'"\1":', temp)
                parsed_plan = json.loads(temp)
            except Exception as e_rec:
                print(f"Advanced recovery failed: {e_rec}. Using regex task block extractor.")
                # Level 3 recovery: Regex parser fallback
                def attempt_regex_plan_extraction(raw_text: str, user_request: str) -> dict:
                    run_title = "Workflow Plan"
                    m_title = re.search(r'"run_title"\s*:\s*"([^"]+)"', raw_text)
                    if m_title:
                        run_title = m_title.group(1)
                    goal = user_request
                    m_goal = re.search(r'"goal"\s*:\s*"([^"]+)"', raw_text)
                    if m_goal:
                        goal = m_goal.group(1)
                    
                    tasks = []
                    # Matches task objects by picking task_id and trailing items up to matching block patterns
                    task_blocks = re.findall(r'\{\s*"task_id".*?\}', raw_text, re.DOTALL)
                    for block in task_blocks:
                        try:
                            tid = re.search(r'"task_id"\s*:\s*"([^"]+)"', block).group(1)
                            title = re.search(r'"title"\s*:\s*"([^"]+)"', block).group(1)
                            desc = re.search(r'"description"\s*:\s*"([^"]+)"', block).group(1)
                            wtype = re.search(r'"worker_type"\s*:\s*"([^"]+)"', block).group(1)
                            
                            dep_match = re.search(r'"depends_on"\s*:\s*\[(.*?)\]', block)
                            depends_on = []
                            if dep_match:
                                depends_on = [x.strip().replace('"', '').replace("'", "") for x in dep_match.group(1).split(",") if x.strip()]
                                
                            tools_match = re.search(r'"allowed_tools"\s*:\s*\[(.*?)\]', block)
                            allowed_tools = []
                            if tools_match:
                                allowed_tools = [x.strip().replace('"', '').replace("'", "") for x in tools_match.group(1).split(",") if x.strip()]
                                
                            success_criteria = "Complete task successfully"
                            sc_match = re.search(r'"success_criteria"\s*:\s*"([^"]+)"', block)
                            if sc_match:
                                success_criteria = sc_match.group(1)
                                
                            tasks.append({
                                "task_id": tid,
                                "title": title,
                                "description": desc,
                                "worker_type": wtype,
                                "depends_on": depends_on,
                                "allowed_tools": allowed_tools,
                                "success_criteria": success_criteria
                            })
                        except Exception:
                            continue
                    
                    if not tasks:
                        raise ValueError("Failed to extract any structured tasks from LLM response")
                    return {"run_title": run_title, "goal": goal, "tasks": tasks}

                parsed_plan = attempt_regex_plan_extraction(raw_text, user_request)

        if validate_plan(parsed_plan):
            return parsed_plan
        else:
            raise ValueError("Failed to validate supervisor workflow plan structure")
    except Exception as e:
        print(f"Supervisor planner failed: {str(e)}")
        raise RuntimeError(f"Workflow Supervisor Planner failed: {str(e)}")


# ── Worker Agent Node Execution ──────────────────────────────────────

async def execute_worker_task(task: TaskDict, prev_tasks_outputs: List[Dict[str, Any]]) -> TaskDict:
    """Executes a single worker agent node using its specific templates and allowed tools."""
    task["status"] = "running"
    task["timestamps"] = task.get("timestamps", {})
    task["timestamps"]["start"] = datetime.datetime.now().isoformat()
    task["retries"] = task.get("retries", 0)
    
    sdir = task.get("storage_dir", "")
    if sdir:
        active_storage_dir.set(sdir)
        
    worker_type = task["worker_type"]
    template = load_templates().get(worker_type)
    if not template:
        task["status"] = "failed"
        task["error"] = f"Unknown worker type: {worker_type}"
        return task
        
    prefs = load_prefs()
    llm = get_atlas_llm(prefs, streaming=False)
    
    # Tool mapping
    tool_map = {
        "web_search": web_search,
        "scrape_page": scrape_page,
        "fetch_url": fetch_url,
        "memory_search": memory_search,
        "artifact_read": artifact_read,
        "file_read": file_read,
        "file_write": file_write,
        "directory_create": directory_create,
        "pdf_export": pdf_export,
        "docx_export": docx_export,
        "pptx_export": pptx_export,
        "excel_export": excel_export,
        "send_email": send_email,
        "browser_automation": browser_automation,
        "shell_exec": shell_exec,
        "escalate_to_supervisor": escalate_to_supervisor
    }
    
    # Resolve allowed tools
    task_tools = [escalate_to_supervisor]  # Always inject supervisor escalation
    allowed_names = set(task.get("allowed_tools", []) + template.get("allowed_tools", []))
    for tool_name in allowed_names:
        if tool_name in tool_map and tool_name != "escalate_to_supervisor":
            task_tools.append(tool_map[tool_name])
            
    # Proactively inject specifically allowed MCP Pool Tools!
    from app.mcp import get_pool_tools
    mcp_defs, tool_sessions, _, _ = get_pool_tools()
    for mdef in mcp_defs:
        mname = mdef["function"]["name"]
        if mname in allowed_names:
            task_tools.append(mdef)
            
    # Inject context from dependencies
    dep_outputs = ""
    for po in prev_tasks_outputs:
        dep_outputs += f"--- Output of Dependent Task {po['task_id']}: {po['title']} ---\n{json.dumps(po.get('output', ''))}\n\n"
        
    sdir = active_storage_dir.get()
    if task.get("worker_type") == "pdf_worker" and sdir:
        # Pre-load all markdown files to avoid agent shortcutting
        md_files = sorted([f for f in os.listdir(sdir) if f.endswith(".md")])
        if md_files:
            dep_outputs += "\n--- PRE-LOADED DETAILED CHAPTER CONTENTS (DO NOT SHORTCUT) ---\n"
            dep_outputs += "Use the following exact file contents to combine and export. DO NOT output brief summaries!\n\n"
            for mf in md_files:
                mfpath = os.path.join(sdir, mf)
                try:
                    with open(mfpath, "r") as f:
                        content = f.read()
                    dep_outputs += f"=== Content of file '{mf}' ===\n{content}\n\n"
                except Exception:
                    pass
        
    system_prompt = template["system_prompt_template"].format(
        task_description=f"{task['description']}\n\nDependency Outputs:\n{dep_outputs}"
    )
    
    correction = task.get("correction_context", "")
    if correction:
        system_prompt += f"\n\n--- CRITICAL RECOVERY INSTRUCTION ---\n{correction}\nAvoid redoing unaffected work if possible."
        
    executed_side_effects = task.get("executed_side_effects", [])
    if executed_side_effects:
        system_prompt += "\n\n--- IDEMPOTENCY WARNING ---\n"
        system_prompt += "You have already executed the following tools successfully in a previous run. DO NOT execute them again to prevent duplicating side-effects:\n"
        for idx, effect in enumerate(executed_side_effects):
            system_prompt += f"{idx+1}. {effect}\n"
        system_prompt += "If the task is complete due to these prior actions, simply return the final output without re-calling the tool.\n"
    
    # ── ReAct Agentic Loop ──────────────────────────────────────────────
    # Each worker iterates: Think → Act (call tool) → Observe (get result)
    # until the LLM stops calling tools (signals task complete) or
    # we hit max_iterations.
    max_iterations = template.get("max_iterations", 20)
    
    retries = 0
    limit = template["retry_limit"]
    
    while retries <= limit:
        try:
            captured_content = ""
            tool_out = ""
            raw_output = ""
            
            # Accumulated step trace for live frontend display
            step_log: list = []
            
            def _append_step(kind: str, content: str, tool_name: str = ""):
                step_log.append({"kind": kind, "tool": tool_name, "content": content[:4000]})
                task["step_log"] = step_log  # persist live for WebSocket streaming
            
            messages = [
                SystemMessage(content=system_prompt),
                HumanMessage(content=(
                    f"Please execute task: {task['title']}.\n"
                    f"Focus on: {task['description']}\n\n"
                    "Work step-by-step. Use your tools one at a time, observe each result, "
                    "and continue until the task is fully complete. "
                    "When you are finished and have all results, produce your final structured output."
                ))
            ]
            
            llm_with_tools = llm.bind_tools(task_tools) if task_tools else llm
            
            # ── Main ReAct Iteration Loop ──────────────────────────────
            for iteration in range(max_iterations):
                resp = await llm_with_tools.ainvoke(messages)
                tool_calls = getattr(resp, "tool_calls", []) or []
                
                if not tool_calls:
                    # ── GUARDRAIL: Enforce mandatory file writes ──
                    tool_names = [getattr(t, "name", "") for t in task_tools] if task_tools else []
                    if "file_write" in tool_names or "file_write" in task.get("allowed_tools", []):
                        if "MUST use the `file_write` tool" in system_prompt and not any(s.get("tool") == "file_write" for s in step_log):
                            # Only nudge if the task actually specifies generating/saving/exporting a file/image
                            desc_lower = (task.get("description", "") + " " + task.get("success_criteria", "")).lower()
                            has_write_instruction = any(keyword in desc_lower for keyword in ["write", "save", "output", "export", "file", "create", "png", "csv", "json", "md", "pdf"])
                            if has_write_instruction:
                                messages.append(AIMessage(content=(resp.content or "").strip()))
                                messages.append(HumanMessage(content="⚠️ GUARDRAIL: You attempted to finish, but you did not use the `file_write` tool to save your output to disk! You MUST save your deliverables using `file_write` before concluding. Please call `file_write` now."))
                                print(f"\n🛡️  [GUARDRAIL] Nudging '{task['title']}' to use file_write before exiting.\n")
                                continue

                    # No tool call → LLM signals it is done
                    raw_output = (resp.content or "").strip()
                    _append_step("final", raw_output)
                    break
                
                # Append AI turn to history
                ai_msg = resp if isinstance(resp, AIMessage) else AIMessage(
                    content=getattr(resp, "content", "") or "",
                    tool_calls=tool_calls
                )
                messages.append(ai_msg)
                
                # Execute each tool call and collect results
                tool_results = []
                for tc in tool_calls:
                    tc_name = tc.get("name", "")
                    tc_args = tc.get("args", {})
                    tc_id   = tc.get("id") or f"call_{iteration}_{tc_name}"
                    
                    _append_step("thought", f"→ {tc_name}({', '.join(f'{k}={str(v)[:60]}' for k,v in tc_args.items())})", tc_name)
                    
                    tool_func = tool_map.get(tc_name)
                    if tool_func:
                        captured_content = (
                            tc_args.get("content", "")
                            or tc_args.get("markdown_content", "")
                            or tc_args.get("query", "")
                            or captured_content
                        )
                        try:
                            single_out = tool_func.invoke(tc_args)
                        except Exception as te:
                            single_out = f"Tool error: {te}"
                    else:
                        # MCP pool tool
                        from app.mcp import get_pool_tools
                        _, tool_sessions, _, _ = get_pool_tools()
                        if tc_name in tool_sessions:
                            session = tool_sessions[tc_name]
                            actual_args = {k: v for k, v in tc_args.items() if k != "ui_status_text"}
                            print(f"[Workflow MCP] '{tc_name}' args={actual_args}")
                            try:
                                result = await asyncio.wait_for(
                                    session.call_tool(tc_name, arguments=actual_args),
                                    timeout=120.0
                                )
                                single_out = "".join([getattr(b, "text", str(b)) for b in result.content])
                            except Exception as me:
                                single_out = f"MCP tool error: {me}"
                        else:
                            single_out = f"Error: Tool '{tc_name}' not found."
                    
                    tool_out = single_out
                    
                    # Side effect tracking (Idempotency)
                    if not str(single_out).startswith("Tool error:") and not str(single_out).startswith("MCP tool error:") and not str(single_out).startswith("Error:"):
                        # Define side effect tools. For now, sending emails or external webhooks
                        if tc_name in ["send_email", "file_write"]:
                            side_effect_record = f"Tool '{tc_name}' called with args: {tc_args}"
                            if "executed_side_effects" not in task:
                                task["executed_side_effects"] = []
                            if side_effect_record not in task["executed_side_effects"]:
                                task["executed_side_effects"].append(side_effect_record)

                    _append_step("observation", str(single_out)[:2000], tc_name)
                    print(f"[ReAct] iter={iteration} tool={tc_name} → {str(single_out)[:150]}")
                    tool_results.append(ToolMessage(content=str(single_out), tool_call_id=tc_id))
                    
                    if str(single_out).startswith("ESCALATION_TRIGGERED:"):
                        raw_output = str(single_out)
                        _append_step("final", raw_output)
                        break
                
                messages.extend(tool_results)
                
                if tool_results and str(tool_results[-1].content).startswith("ESCALATION_TRIGGERED:"):
                    break
                
                # On the penultimate iteration, nudge toward final output
                if iteration == max_iterations - 2:
                    messages.append(HumanMessage(
                        content="You have used many steps. Now produce your final answer without calling any more tools."
                    ))
            else:
                # Exhausted iterations — force a final answer
                resp = await llm.ainvoke(messages + [
                    HumanMessage(content="Produce your final structured output based on everything done so far.")
                ])
                raw_output = (resp.content or "").strip()
                _append_step("final", raw_output)
                
            # Parse output
            try:
                if raw_output.startswith("ESCALATION_TRIGGERED:"):
                    task["error"] = raw_output.replace("ESCALATION_TRIGGERED:", "").strip()
                    print(f"\n⚠️  [ESCALATION] Agent '{task['title']}' called Supervisor: {task['error']}\n")
                    raise Exception(f"Agent explicitly escalated: {task['error']}")
                # Strip <think> tags from reasoning models (DeepSeek R1, etc)
                clean_text = raw_output
                if "<think>" in clean_text:
                    clean_text = _re.sub(r'<think>.*?</think>', '', clean_text, flags=_re.DOTALL).strip()
                # Strip markdown tags if LLM wrapped in JSON block
                if "```json" in clean_text:
                    clean_text = clean_text.split("```json")[1].split("```")[0].strip()
                elif "```" in clean_text:
                    clean_text = clean_text.split("```")[1].split("```")[0].strip()
                
                # Robust JSON single-quote key and value repair prior to loads
                import re as _re_repair
                try:
                    output_val = json.loads(clean_text)
                except Exception:
                    repaired = clean_text
                    repaired = _re_repair.sub(r"'\s*:\s*", '": ', repaired)
                    repaired = _re_repair.sub(r"{\s*'", '{"', repaired)
                    repaired = _re_repair.sub(r",\s*'", ',"', repaired)
                    repaired = _re_repair.sub(r"'\s*,", '",', repaired)
                    repaired = _re_repair.sub(r"'\s*}", '"}', repaired)
                    repaired = _re_repair.sub(r"'\s*]", '"]', repaired)
                    repaired = _re_repair.sub(r"\[\s*'", '["', repaired)
                    repaired = _re_repair.sub(r',\s*([\]}])', r'\1', repaired)
                    output_val = json.loads(repaired)
            except Exception:
                output_val = {"result": raw_output}
                
            # Recovery fail-safe from captured tool content if output is empty or truncated
            if isinstance(output_val, dict):
                res = output_val.get("result") or output_val.get("markdown") or output_val.get("summary") or output_val.get("content") or output_val.get("findings") or ""
                has_substance = str(res).strip() and len(str(res).strip()) >= 10
                
                # 1. Recover from tool execution result
                if not has_substance and 'tool_out' in locals() and tool_out and len(str(tool_out).strip()) > 10:
                    output_val["result"] = str(tool_out)
                    output_val["summary"] = "Synthesized automatically from execution tool results."
                    has_substance = True
                    
                # 2. Recover from captured tool call arguments (content/markdown_content)
                if not has_substance and captured_content and len(captured_content.strip()) > 10:
                    output_val["markdown"] = captured_content
                    has_substance = True
                elif not output_val.get("markdown") and captured_content and len(captured_content.strip()) > 10:
                    output_val["markdown"] = captured_content
                    
                # 3. Last resort: recover from raw LLM output if it has substance but JSON parse stripped it  
                if not has_substance and raw_output and len(raw_output.strip()) > 50:
                    output_val["result"] = raw_output
                    output_val["summary"] = "Recovered from raw LLM response."
                    
            # Fail-safe writer folder sync to ensure chapters are ALWAYS written to files!
            if task.get("worker_type") == "writer" and sdir:
                import re as _re_fs
                match = _re_fs.search(r'[\'"]?(ch\d+\.md)[\'"]?', task.get("description", ""))
                if not match:
                    match = _re_fs.search(r'[\'"]?(ch\d+\.md)[\'"]?', task.get("title", ""))
                if match:
                    filename = match.group(1)
                    filepath = os.path.join(sdir, filename)
                    
                    file_written_ok = False
                    # Check if file was already successfully written to disk by the tool call
                    if os.path.exists(filepath) and os.path.getsize(filepath) > 100:
                        try:
                            with open(filepath, "r") as f:
                                existing_content = f.read()
                            cleaned_existing = clean_markdown_if_json(existing_content)
                            if cleaned_existing != existing_content:
                                with open(filepath, "w") as f:
                                    f.write(cleaned_existing)
                                print(f"Cleaned JSON wrapper from existing on-disk file: {filepath}")
                        except Exception as e:
                            print(f"Error checking on-disk JSON wrapper: {str(e)}")
                        file_written_ok = True
                        print(f"Verified chapter file exists on disk: {filepath} ({os.path.getsize(filepath)} bytes)")
                        
                    if not file_written_ok:
                        markdown_content = ""
                        if isinstance(output_val, dict):
                            markdown_content = output_val.get("markdown") or output_val.get("result") or ""
                        elif isinstance(output_val, str):
                            markdown_content = output_val
                            
                        # Recover from captured first-turn tool call if JSON output got truncated/escaped!
                        if (not markdown_content or len(markdown_content.strip()) < 100) and captured_content:
                            markdown_content = captured_content
                            
                        markdown_content = clean_markdown_if_json(markdown_content)
                        if markdown_content and len(markdown_content) > 100:
                            try:
                                with open(filepath, "w") as f:
                                    f.write(markdown_content)
                                print(f"Fail-safe wrote writer chapter to {filepath}")
                                file_written_ok = True
                            except Exception as e:
                                print(f"Fail-safe writer failed: {str(e)}")
                                
                    # If empty or not written, throw Value error to trigger exponential backoff retry!
                    if not file_written_ok:
                        raise ValueError(f"Technical Writer received empty or throttled response for {filename}. Retrying...")
                            
            task["output"] = output_val
            task["status"] = "completed"
            task["timestamps"]["end"] = datetime.datetime.now().isoformat()
            
            # Record artifacts
            artifacts = []
            
            # For pdf_worker: if LLM didn't actually call pdf_export, do it ourselves
            if worker_type == "pdf_worker":
                already_done = False
                if isinstance(output_val, dict) and output_val.get("pdf_path") and "Successfully" in str(output_val.get("pdf_path")):
                    already_done = True
                    
                if not already_done:
                    md_content = ""
                    parts_count = 0
                    # 1. Try reading actual written markdown files from the workspace directory first!
                    if sdir:
                        md_files = sorted([f for f in os.listdir(sdir) if f.endswith(".md")])
                        if md_files:
                            parts = []
                            for mf in md_files:
                                try:
                                    with open(os.path.join(sdir, mf), "r") as f:
                                        parts.append(clean_markdown_if_json(f.read()))
                                except Exception:
                                    pass
                            md_content = "\n\n\\pagebreak\n\n".join(parts)
                            parts_count = len(md_files)
                            
                    # 2. Fallback: join all dependent task outputs sequentially!
                    if not md_content:
                        parts = []
                        for po in prev_tasks_outputs:
                            out = po.get("output", {})
                            part = ""
                            if isinstance(out, dict):
                                part = out.get("markdown") or out.get("result") or ""
                            elif isinstance(out, str):
                                part = out
                            part = clean_markdown_if_json(part)
                            if part and len(part.strip()) > 50:
                                parts.append(part)
                        md_content = "\n\n\\pagebreak\n\n".join(parts)
                        parts_count = len(parts)
                        
                    if md_content:
                        # 1. Determine dynamic safe name
                        wf_title = active_workflow_title.get()
                        if wf_title:
                            safe_name = _re.sub(r'[^a-zA-Z0-9_\-]', '_', wf_title).strip("_")[:50]
                        else:
                            safe_name = _re.sub(r'[^a-zA-Z0-9_\-]', '_', task.get("title", "report")).strip("_")[:40]

                        # 2. Determine which formats are requested dynamically
                        title_lower = task.get("title", "").lower()
                        desc_lower = task.get("description", "").lower()
                        
                        has_pdf_req = "pdf" in title_lower or "pdf" in desc_lower
                        has_docx_req = "docx" in title_lower or "word" in title_lower or "docx" in desc_lower or "word" in desc_lower
                        has_pptx_req = "pptx" in title_lower or "powerpoint" in title_lower or "pptx" in desc_lower or "powerpoint" in desc_lower or "slide" in title_lower or "slide" in desc_lower
                        has_xlsx_req = "excel" in title_lower or "xlsx" in title_lower or "excel" in desc_lower or "xlsx" in desc_lower or "spreadsheet" in title_lower or "spreadsheet" in desc_lower
                        
                        compile_pdf = True
                        compile_docx = True
                        compile_pptx = True
                        compile_xlsx = True
                        
                        if has_pdf_req or has_docx_req or has_pptx_req or has_xlsx_req:
                            compile_pdf = has_pdf_req
                            compile_docx = has_docx_req
                            compile_pptx = has_pptx_req
                            compile_xlsx = has_xlsx_req
                            
                        # 2. Render premium PDF booklet
                        if compile_pdf:
                            pdf_result = pdf_export.invoke({"markdown_content": md_content, "filename": f"{safe_name}.pdf"})
                            print(f"[PDF Worker] Direct PDF invocation result: {pdf_result}")
                            
                            if "Successfully" in pdf_result:
                                target_path = pdf_result.split("at: ")[1].strip()
                                output_val = {
                                    "status": "success",
                                    "pdf_path": target_path,
                                    "result": f"Successfully compiled document booklet containing {parts_count} chapters."
                                }
                                path_match = _re.search(r'at:\s*(.+)$', pdf_result)
                                if path_match:
                                    artifacts.append(path_match.group(1).strip())
                                task["output"] = output_val
                            
                        # 3. Render premium editable Word Document (.docx)
                        if compile_docx:
                            try:
                                docx_res = docx_export.invoke({"markdown_content": md_content, "filename": f"{safe_name}.docx"})
                                print(f"[PDF Worker] Word DOCX generation result: {docx_res}")
                                if "Successfully" in docx_res:
                                    path_match = _re.search(r'at:\s*(.+)$', docx_res)
                                    if path_match:
                                        artifacts.append(path_match.group(1).strip())
                            except Exception as de:
                                print(f"[PDF Worker] Word document export failed: {de}")
                        # 4. Render premium slide deck PowerPoint (.pptx)
                        if compile_pptx:
                            try:
                                pptx_res = pptx_export.invoke({"markdown_content": md_content, "filename": f"{safe_name}.pptx"})
                                print(f"[PDF Worker] PowerPoint PPTX generation result: {pptx_res}")
                                if "Successfully" in pptx_res:
                                    path_match = _re.search(r'at:\s*(.+)$', pptx_res)
                                    if path_match:
                                        artifacts.append(path_match.group(1).strip())
                            except Exception as pe:
                                print(f"[PDF Worker] PowerPoint presentation export failed: {pe}")
    
                        # 5. Render styled Excel spreadsheet (.xlsx) ONLY if tabular data is present in the markdown report!
                        if compile_xlsx and "|" in md_content:
                            try:
                                tables_data = []
                                current_table = []
                                for line in md_content.split("\n"):
                                    line_strip = line.strip()
                                    if line_strip.startswith("|") and line_strip.endswith("|"):
                                        parts = [p.strip() for p in line_strip.split("|")[1:-1]]
                                        # Skip separators
                                        if all(set(p).issubset({'-', ':', ' '}) for p in parts):
                                            continue
                                        current_table.append(parts)
                                    else:
                                        if current_table:
                                            tables_data.append(current_table)
                                            current_table = []
                                if current_table:
                                    tables_data.append(current_table)
                                    
                                if tables_data:
                                    xls_res = excel_export.invoke({
                                        "table_data_json": json.dumps(tables_data[0]),
                                        "filename": f"{safe_name}.xlsx"
                                    })
                                    print(f"[PDF Worker] Excel XLSX generation result: {xls_res}")
                                    if "Successfully" in xls_res:
                                        path_match = _re.search(r'at:\s*(.+)$', xls_res)
                                        if path_match:
                                            artifacts.append(path_match.group(1).strip())
                            except Exception as ee:
                                print(f"[PDF Worker] Excel spreadsheet export failed: {ee}")
            
            # Standard artifact detection
            if not artifacts:
                if isinstance(output_val, dict) and "pdf_path" in output_val:
                    p = output_val["pdf_path"]
                    # Could be a result string like "Successfully rendered PDF at: /path"
                    path_match = _re.search(r'at:\s*(.+?)(?:\s*$)', p)
                    if path_match:
                        artifacts.append(path_match.group(1).strip())
                    elif os.path.exists(p):
                        artifacts.append(p)
                elif isinstance(output_val, dict) and "file_path" in output_val:
                    artifacts.append(output_val["file_path"])
            
            # Fallback: scan uploads dir for any recently created PDFs (within last 60s)
            if not artifacts and worker_type == "pdf_worker":
                uploads_dir = os.path.join(DATA_DIR, "uploads")
                if os.path.isdir(uploads_dir):
                    now_ts = time.time()
                    for fname in os.listdir(uploads_dir):
                        if fname.endswith(".pdf"):
                            fpath = os.path.join(uploads_dir, fname)
                            if now_ts - os.path.getmtime(fpath) < 60:
                                artifacts.append(fpath)
            
            task["artifacts"] = artifacts
            
            # --- Content Validation Layer ---
            # Checks if the output actually contains real substance before marking complete.
            validation_passed = True
            validation_error = ""
            
            # Extract content to validate
            substance_text = ""
            if isinstance(output_val, dict):
                substance_text = output_val.get("markdown") or output_val.get("result") or output_val.get("summary") or output_val.get("content") or output_val.get("notes") or output_val.get("findings") or ""
                if not substance_text:
                    str_vals = [str(v) for v in output_val.values() if isinstance(v, str) and len(str(v).strip()) > 5]
                    substance_text = "\n".join(str_vals)
            elif isinstance(output_val, str):
                substance_text = output_val
                
            substance_text = substance_text.strip()
            
            # Pre-validation recovery: if substance is empty/short, attempt to rescue before rejecting
            if len(substance_text) < 30:
                if 'tool_out' in locals() and tool_out and len(str(tool_out).strip()) > 30:
                    substance_text = str(tool_out).strip()
                    output_val["result"] = substance_text
                    task["output"] = output_val
                elif captured_content and len(captured_content.strip()) > 30:
                    substance_text = captured_content.strip()
                    output_val["markdown"] = substance_text
                    task["output"] = output_val
                elif raw_output and len(raw_output.strip()) > 50:
                    substance_text = raw_output.strip()
                    output_val["result"] = substance_text
                    task["output"] = output_val
            
            # Obvious placeholder / failure patterns
            failure_patterns = [
                "failed to execute", 
                "could not retrieve", 
                "throttled", 
                "rate limit", 
                "i apologize, but", 
                "i cannot complete",
                "internal server error",
                "failed with status 500",
                "placeholder content"
            ]
            
            # Apply checks based on worker type
            if worker_type in ("researcher", "writer"):
                if len(substance_text) < 100:
                    validation_passed = False
                    validation_error = f"Content is too short ({len(substance_text)} chars, expected >= 100). Substance check failed."
                elif any(pat in substance_text.lower() for pat in failure_patterns):
                    validation_passed = False
                    validation_error = "Content contains failure placeholders or error messages. Substance check failed."
            elif worker_type == "pdf_worker":
                # Ensure at least one PDF artifact was created and it exists on disk with size > 1024 bytes
                has_pdf = False
                for art in artifacts:
                    if art.endswith(".pdf") and os.path.exists(art) and os.path.getsize(art) > 1024:
                        has_pdf = True
                        break
                # Fallback: check storage dir for recently created PDFs
                if not has_pdf and sdir and os.path.isdir(sdir):
                    for f in os.listdir(sdir):
                        fpath = os.path.join(sdir, f)
                        if f.endswith(".pdf") and os.path.exists(fpath) and os.path.getsize(fpath) > 1024:
                            artifacts.append(fpath)
                            task["artifacts"] = artifacts
                            has_pdf = True
                            break
                if not has_pdf:
                    validation_passed = False
                    validation_error = "Generated PDF file is missing, empty, or corrupted. Substance check failed."
            else:
                # Generic fallback check for all other agents
                if len(substance_text) < 20:
                    validation_passed = False
                    validation_error = f"Output content is too short ({len(substance_text)} chars, expected >= 20). Substance check failed."
                elif any(pat in substance_text.lower() for pat in failure_patterns):
                    validation_passed = False
                    validation_error = "Output contains obvious failure patterns."
                    
            if not validation_passed:
                print(f"[Validation Layer] Rejecting task '{task['title']}' output. Error: {validation_error}")
                raise ValueError(f"Task output failed substance validation: {validation_error}")

            return task
        except Exception as e:
            retries += 1
            task["retries"] = retries
            task["error"] = str(e)
            await asyncio.sleep(2 ** retries) # Exponential Backoff
            
    task["status"] = "failed"
    return task


# ── LangGraph Workflow State Orchestrator ────────────────────────────

def run_scheduler_step(state: WorkflowState) -> WorkflowState:
    """Represent scheduler executions inside a LangGraph graph execution node."""
    # Find all ready tasks
    completed = set(state["completed_tasks"])
    failed = set(state["failed_tasks"])
    running = set(state["running_tasks"])
    
    for task in state["tasks"]:
        tid = task["task_id"]
        if tid in completed or tid in failed or tid in running:
            continue
            
        # Check backoff for recovery retries
        backoff_until = task.get("retry_backoff_until")
        if backoff_until and datetime.datetime.now().timestamp() < backoff_until:
            continue
            
        # Check dependencies
        deps = set(task.get("depends_on", []))
        if deps.issubset(completed):
            # Ready to schedule!
            state["ready_queue"].append(tid)
            
    return state

class MultiAgentWorkflowEngine:
    """Core multi-agent lifecycle coordinator orchestrating LangGraph plans."""
    active_tasks = {}  # run_id -> asyncio.Task
    
    @staticmethod
    async def run_supervisor_recovery(state: WorkflowState, failed_task: TaskDict) -> TaskDict:
        prefs = load_prefs()
        llm = get_atlas_llm(prefs, streaming=False)
        
        # Circuit Breaker Logic
        current_error = failed_task.get('error', 'Unknown')
        error_signature = current_error[:100]  # simple signature based on first 100 chars
        
        if failed_task.get('last_error_signature') == error_signature:
            failed_task['error_signature_count'] = failed_task.get('error_signature_count', 0) + 1
        else:
            failed_task['last_error_signature'] = error_signature
            failed_task['error_signature_count'] = 1
            
        if failed_task['error_signature_count'] >= 3:
            failed_task["failure_type"] = "repeated_failure"
            failed_task["recovery_decision"] = "fail_terminal"
            failed_task["correction_context"] = "Circuit breaker triggered: Same error signature occurred 3 times."
            return failed_task

        prompt = f"""You are the Workflow Recovery Supervisor.
A worker node has failed its execution. Your job is to classify the failure and decide the safest recovery action.
If the failure is caused by bad reasoning, malformed output, or validation errors, you MUST inject structured feedback so it corrects itself on the next run.

FAILURE EVENT REPORT:
- Node ID: {failed_task['task_id']}
- Title: {failed_task['title']}
- Description: {failed_task['description']}
- Attempt Count: {failed_task.get('attempt_count', 1)}
- Last Error Message: {current_error}
- Step Traces (Last 3): {json.dumps(failed_task.get('step_log', [])[-3:]) if failed_task.get('step_log') else '[]'}

CLASSIFICATION CATEGORIES (Choose one):
- transient_error
- timeout_network
- malformed_output
- validation_failure
- dependency_failure
- permission_auth_failure
- unrecoverable

RECOVERY ACTIONS (Choose one):
- retry_node (Simple retry for transient issues)
- retry_with_context (Retry and provide correction_context for logic/validation fixes)
- restart_subtree (Restart this node and all descendants, preserving unaffected siblings)
- replan_workflow (Invalidate the DAG and generate a completely new plan)
- full_restart (Restart the entire workflow only if root assumptions are broken)
- fail_terminal (Unrecoverable, stop workflow)

Respond in JSON exactly matching this schema:
{{
  "classification": "string",
  "action": "string",
  "correction_context": "string (If action is retry_with_context, provide clear instructions on how to fix the error. Otherwise empty.)"
}}
"""
        import re
        try:
            resp = await llm.ainvoke([SystemMessage(content=prompt)])
            text = (resp.content or "").strip()
            if "```json" in text:
                text = text.split("```json")[1].split("```")[0].strip()
            elif "```" in text:
                text = text.split("```")[1].split("```")[0].strip()
            decision = json.loads(text)
        except Exception as e:
            decision = {
                "classification": "unrecoverable",
                "action": "fail_terminal",
                "correction_context": f"Supervisor failed to parse recovery: {str(e)}"
            }
            
        failed_task["failure_type"] = decision.get("classification", "unknown")
        failed_task["recovery_decision"] = decision.get("action", "fail_terminal")
        failed_task["correction_context"] = decision.get("correction_context", "")
        
        print(f"\n👨‍⚕️ [SUPERVISOR] Evaluated Task '{failed_task['title']}'")
        print(f"   ↳ Classification: {failed_task['failure_type']}")
        print(f"   ↳ Action: {failed_task['recovery_decision']}")
        if failed_task["correction_context"]:
            print(f"   ↳ Context: {failed_task['correction_context'][:150]}...")
            
        return failed_task

    @staticmethod
    async def create_run(user_request: str, user_id: str = "admin") -> str:
        run_id = f"wf_{uuid.uuid4().hex[:8]}"
        
        # Create per-workflow storage directory
        wf_dir = os.path.join(DATA_DIR, "workflows", run_id)
        os.makedirs(wf_dir, exist_ok=True)
        
        plan = await run_supervisor_planner(user_request)
        
        tasks_list: List[TaskDict] = []
        for t in plan["tasks"]:
            tasks_list.append({
                "task_id": t["task_id"],
                "title": t["title"],
                "description": t["description"],
                "worker_type": t["worker_type"],
                "depends_on": t.get("depends_on", []),
                "allowed_tools": t["allowed_tools"],
                "success_criteria": t.get("success_criteria", ""),
                "expected_output_schema": t.get("expected_output_schema", {}),
                "status": "pending",
                "retries": 0,
                "artifacts": []
            })
            
        # Determine a cool icon based on title/goal
        icon = "fa-circle-nodes"
        lower_title = plan.get("run_title", "").lower()
        lower_goal = plan.get("goal", "").lower()
        combined = lower_title + " " + lower_goal
        if "pdf" in combined or "report" in combined:
            icon = "fa-file-pdf"
        elif "research" in combined or "study" in combined or "virus" in combined or "hanta" in combined or "microcomputer" in combined:
            icon = "fa-magnifying-glass"
        elif "code" in combined or "compile" in combined or "program" in combined:
            icon = "fa-code"
        elif "email" in combined or "mail" in combined:
            icon = "fa-envelope"
        elif "style" in combined or "design" in combined:
            icon = "fa-wand-magic-sparkles"
            
        wf_state = {
            "run_id": run_id,
            "user_id": user_id,
            "user_request": user_request,
            "run_title": plan["run_title"],
            "run_icon": icon,
            "plan": plan,
            "tasks": tasks_list,
            "ready_queue": [],
            "running_tasks": [],
            "completed_tasks": [],
            "failed_tasks": [],
            "artifacts": [],
            "notifications": [f"💬 User: {user_request}"],
            "collaborative_chat": [{
                "role": "user",
                "sender": "User",
                "message": user_request,
                "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            }],
            "debug_logs": [],
            "final_result": None,
            "status": "pending",
            "created_at": datetime.datetime.now().isoformat(),
            "storage_dir": wf_dir
        }
        
        # Save to database
        db = load_workflows()
        db[run_id] = wf_state
        save_workflows(db)
        return run_id

    @staticmethod
    async def execute_run(run_id: str):
        import asyncio
        MultiAgentWorkflowEngine.active_tasks[run_id] = asyncio.current_task()
        try:
            await MultiAgentWorkflowEngine._execute_run_internal(run_id)
        finally:
            MultiAgentWorkflowEngine.active_tasks.pop(run_id, None)

    @staticmethod
    async def _execute_run_internal(run_id: str):
        db = load_workflows()
        if run_id not in db:
            return

        def _dbg(state, msg):
            ts = datetime.datetime.now().strftime("%H:%M:%S.%f")[:-3]
            state.setdefault("debug_logs", []).append(f"[{ts}] {msg}")

        state = db[run_id]
        state["status"] = "running"
        state["started_at"] = datetime.datetime.now().isoformat()
        active_workflow_title.set(state.get("run_title", ""))
        active_storage_dir.set(state.get("storage_dir", ""))
        _dbg(state, f"🚀 Workflow '{state['run_title']}' execution started")
        _dbg(state, f"📋 Plan contains {len(state['tasks'])} task nodes")
        state["notifications"].append(f"Execution started for '{state['run_title']}'")
        db[run_id] = state
        save_workflows(db)
        
        while state["status"] == "running":
            # Check for external manual cancellation/failure
            db_check = load_workflows()
            if db_check.get(run_id, {}).get("status") in ["failed", "completed", "cancelled"]:
                _dbg(state, "🛑 Workflow loop halted due to status change")
                break
                
            state = run_scheduler_step(state)
            
            ready = list(state["ready_queue"])
            if not ready and not state["running_tasks"]:
                _dbg(state, "✅ All task nodes processed — exiting main loop")
                break
                
            if not ready:
                await asyncio.sleep(1)
                db = load_workflows()
                state = db[run_id]
                continue
                
            state["ready_queue"] = []
            
            spawn_tasks = []
            for tid in ready:
                state["running_tasks"].append(tid)
                target_task = None
                for t in state["tasks"]:
                    if t["task_id"] == tid:
                        target_task = t
                        break

                _dbg(state, f"⚡ Spawning worker '{target_task['worker_type']}' for task '{target_task['title']}' ({tid})")
                state["notifications"].append(f"Task '{target_task['title']}' started ({target_task['worker_type']})")
                        
                dep_ids = target_task.get("depends_on", [])
                dep_outputs = [dt for dt in state["tasks"] if dt["task_id"] in dep_ids]
                        
                target_task["storage_dir"] = state.get("storage_dir", "")
                target_task["status"] = "running"
                target_task["timestamps"] = target_task.get("timestamps", {})
                target_task["timestamps"]["start"] = datetime.datetime.now().isoformat()
                
                spawn_tasks.append(execute_worker_task(target_task, dep_outputs))
                
            db[run_id] = state
            save_workflows(db)
            
            completed_nodes = await asyncio.gather(*spawn_tasks)
            
            db = load_workflows()
            state = db[run_id]
            
            for cn in completed_nodes:
                tid = cn["task_id"]
                for idx, t in enumerate(state["tasks"]):
                    if t["task_id"] == tid:
                        state["tasks"][idx] = cn
                        break
                        
                if tid in state["running_tasks"]:
                    state["running_tasks"].remove(tid)
                if cn["status"] == "completed":
                    state["completed_tasks"].append(tid)
                    _dbg(state, f"✅ Task '{cn['title']}' completed successfully (retries: {cn.get('retries', 0)})")
                    state["notifications"].append(f"Task '{cn['title']}' completed ✓")
                    if cn.get("artifacts"):
                        state["artifacts"].extend(cn["artifacts"])
                        _dbg(state, f"📦 Artifacts produced: {cn['artifacts']}")
                        # Copy generated artifacts to workflow's isolated storage folder
                        sdir = state.get("storage_dir")
                        if sdir:
                            import shutil
                            os.makedirs(sdir, exist_ok=True)
                            for art in cn["artifacts"]:
                                if os.path.exists(art):
                                    try:
                                        shutil.copy2(art, sdir)
                                        _dbg(state, f"💾 Saved artifact {os.path.basename(art)} to workflow storage directory")
                                    except Exception as ce:
                                        _dbg(state, f"⚠️ Failed to copy artifact to storage: {str(ce)}")
                else:
                    _dbg(state, f"🚨 Task '{cn['title']}' failed. Escalating to Supervisor for recovery evaluation...")
                    
                    cn["attempt_count"] = cn.get("attempt_count", 0) + 1
                    
                    cn = await MultiAgentWorkflowEngine.run_supervisor_recovery(state, cn)
                    decision = cn.get("recovery_decision", "fail_terminal")
                    classification = cn.get("failure_type", "unknown")
                    _dbg(state, f"🏥 Supervisor Decision: {decision} (Class: {classification})")
                    
                    if decision in ["retry_node", "retry_with_context"] and cn["attempt_count"] <= 3:
                        # Re-queue the task
                        cn["status"] = "pending"
                        if decision == "retry_with_context":
                            _dbg(state, "💉 Injecting correction context for next attempt.")
                        
                        # Apply backoff
                        backoff = 2 ** cn["attempt_count"]
                        cn["retry_backoff_until"] = datetime.datetime.now().timestamp() + backoff
                        
                        state["notifications"].append(f"Task '{cn['title']}' scheduled for retry ({cn['attempt_count']}/3) ↻")
                    elif decision == "restart_subtree":
                        _dbg(state, f"🔄 Subtree restart initiated from node '{tid}'. Unaffected completed nodes will be preserved.")
                        cn["status"] = "pending"
                        cn["attempt_count"] = 0
                        
                        # Recursively find descendants to invalidate
                        descendants = set([tid])
                        changed = True
                        while changed:
                            changed = False
                            for t in state["tasks"]:
                                if t["task_id"] not in descendants and any(dep in descendants for dep in t.get("depends_on", [])):
                                    descendants.add(t["task_id"])
                                    changed = True
                                    
                        for t in state["tasks"]:
                            if t["task_id"] in descendants and t["task_id"] != tid:
                                t["status"] = "pending"
                                t["attempt_count"] = 0
                                if t["task_id"] in state["completed_tasks"]:
                                    state["completed_tasks"].remove(t["task_id"])
                                    
                        state["notifications"].append(f"Subtree restarted from '{cn['title']}' 🔄")
                        
                    elif decision == "replan_workflow":
                        _dbg(state, f"🧠 Workflow DAG Replanning initiated! (Failing out for manual user adjustment)")
                        cn["status"] = f"failed_{classification}"
                        state["failed_tasks"].append(tid)
                        state["notifications"].append(f"Workflow needs replanning. Terminal failure ✗")
                        
                    elif decision == "full_restart":
                        _dbg(state, f"⚠️ Full workflow restart initiated! Root assumptions were broken.")
                        for t in state["tasks"]:
                            t["status"] = "pending"
                            t["attempt_count"] = 0
                        state["completed_tasks"] = []
                        state["failed_tasks"] = []
                        state["running_tasks"] = []
                        state["notifications"].append("Entire workflow restarted ⚠️")
                        
                    else:
                        cn["status"] = f"failed_{classification}"
                        state["failed_tasks"].append(tid)
                        _dbg(state, f"❌ Task '{cn['title']}' FAILED TERMINALLY: {cn.get('error', 'unknown')}")
                        state["notifications"].append(f"Task '{cn['title']}' failed unrecoverably ✗")
                    
            db[run_id] = state
            save_workflows(db)
            
        # Finalization
        db = load_workflows()
        state = db[run_id]
        
        state["completed_at"] = datetime.datetime.now().isoformat()
        
        if state["failed_tasks"]:
            state["status"] = "failed"
            state["final_result"] = f"Workflow failed. Errored tasks: {', '.join(state['failed_tasks'])}."
            _dbg(state, f"🔴 Workflow FAILED — {len(state['failed_tasks'])} task(s) errored")
            
            # Let's generate a failures final response
            try:
                prefs = load_prefs()
                llm = get_atlas_llm(prefs, streaming=False)
                prompt = (
                    "You are the Top-Level Workflow Supervisor. The workflow execution has encountered terminal errors and stopped.\n"
                    f"Objective/Goal: {state['user_request']}\n"
                    f"Failed Tasks: {', '.join(state['failed_tasks'])}\n\n"
                    "Please write a professional, polite response to the user. "
                    "Explain that the workflow encountered an error and could not complete all tasks. "
                    "List the errored tasks and express that they can ask you to adjust parameters or retry."
                )
                response = await llm.ainvoke([HumanMessage(content=prompt)])
                final_text = response.content.strip()
                
                state["notifications"].append(f"🤖 Supervisor: {final_text}")
                state.setdefault("collaborative_chat", []).append({
                    "role": "assistant",
                    "sender": "Supervisor",
                    "message": final_text,
                    "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                })
                state["final_result"] = final_text
            except Exception as e_sup:
                print(f"Error generating supervisor failure response: {e_sup}")
        else:
            state["status"] = "completed"
            state["final_result"] = "All tasks completed successfully. Artifacts ready for download."
            _dbg(state, f"🟢 Workflow COMPLETED — {len(state['completed_tasks'])} task(s) succeeded")
            
            # Let's call the supervisor to generate a final completed message summarizing everything and presenting the artifacts!
            try:
                import os
                # Look for artifacts and files in state
                storage_dir = state.get("storage_dir", "")
                generated_files = []
                if storage_dir and os.path.exists(storage_dir):
                    for f in os.listdir(storage_dir):
                        if os.path.isfile(os.path.join(storage_dir, f)) and not f.startswith("."):
                            generated_files.append(f)
                            
                # Populate artifacts list so the frontend can find them
                state["artifacts"] = generated_files
                
                tasks_summary = ""
                for t in state.get("tasks", []):
                    tasks_summary += f"- {t['title']}: {t.get('status', 'pending')}\n"
                    
                prefs = load_prefs()
                llm = get_atlas_llm(prefs, streaming=False)
                prompt = (
                    "You are the Top-Level Workflow Supervisor. The workflow execution has finished successfully!\n"
                    f"Objective/Goal: {state['user_request']}\n"
                    f"Tasks Executed:\n{tasks_summary}\n"
                    f"Generated Files: {', '.join(generated_files) if generated_files else 'None'}\n\n"
                    "Please write a highly polished, professional, final response. "
                    "Confirm the completion of the objective. "
                    "Provide a concise, high-value summary of what was accomplished and the key findings. "
                    "Be polite, expert, and premium."
                )
                response = await llm.ainvoke([HumanMessage(content=prompt)])
                final_text = response.content.strip()
                
                # Append to collaborative chat and notifications
                state["notifications"].append(f"🤖 Supervisor: {final_text}")
                state.setdefault("collaborative_chat", []).append({
                    "role": "assistant",
                    "sender": "Supervisor",
                    "message": final_text,
                    "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                })
                state["final_result"] = final_text
            except Exception as e_sup:
                print(f"Error generating supervisor final summary: {e_sup}")
            
        state["notifications"].append(f"Workflow '{state['run_title']}' finished: {state['status'].upper()}")
        db[run_id] = state
        save_workflows(db)

